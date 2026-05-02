"""PowerPoint Exporter for Analytics Reports

Generates professional PowerPoint presentations from analytics data with charts,
tables, and professional formatting. Supports python-pptx (full features) with
graceful fallback to PDF export.

Installation:
    Preferred (full features):
        pip install python-pptx Pillow

    Fallback:
        Uses PDFExporter (requires reportlab)

Features:
    - Professional 16:9 presentation layout
    - Title slide with metadata
    - Table of contents
    - Section slides with metrics and charts
    - Embedded chart images
    - Consistent branding and styling
    - Contact/summary slide

Example:
    >>> from analytics.exporters import PPTXExporter
    >>> exporter = PPTXExporter()
    >>> report_data = {
    ...     "metadata": {
    ...         "generated_at": "2026-05-01T12:00:00Z",
    ...         "report_type": "summary",
    ...         "date_range": {"start": "2026-04-01", "end": "2026-04-30"}
    ...     },
    ...     "sections": [
    ...         {
    ...             "title": "Performance Metrics",
    ...             "metrics": {"total_users": 15234, "revenue": 125430.50},
    ...             "charts": []
    ...         }
    ...     ]
    ... }
    >>> success, result = exporter.export_to_pptx(report_data, "report.pptx")
    >>> if success:
    ...     print(f"Presentation saved to: {result}")
"""

import os
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Tuple, Optional

# Try importing python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

# Try importing Pillow for image validation
try:
    from PIL import Image as PILImage
    HAS_PILLOW = True
except ImportError:
    HAS_PILLOW = False


# Color scheme for professional presentations
COLORS = {
    'primary': RGBColor(44, 62, 80) if HAS_PYTHON_PPTX else '#2c3e50',      # Dark blue-gray
    'secondary': RGBColor(52, 152, 219) if HAS_PYTHON_PPTX else '#3498db',  # Blue
    'accent_red': RGBColor(231, 76, 60) if HAS_PYTHON_PPTX else '#e74c3c',  # Red
    'accent_green': RGBColor(46, 204, 113) if HAS_PYTHON_PPTX else '#2ecc71', # Green
    'light_gray': RGBColor(236, 240, 241) if HAS_PYTHON_PPTX else '#ecf0f1', # Light gray
    'dark_gray': RGBColor(127, 140, 141) if HAS_PYTHON_PPTX else '#7f8c8d',  # Dark gray
    'white': RGBColor(255, 255, 255) if HAS_PYTHON_PPTX else '#ffffff',      # White
}


class PPTXExporter:
    """Professional PowerPoint presentation generator with graceful fallback"""

    def __init__(self):
        """Initialize PowerPoint exporter with available libraries"""
        self.has_python_pptx = HAS_PYTHON_PPTX
        self.has_pillow = HAS_PILLOW

        if not self.has_python_pptx:
            print("WARNING: python-pptx not available - PowerPoint exports will use PDF fallback")
            print("   Install with: pip install python-pptx Pillow")

    def export_to_pptx(
        self,
        report_data: Dict[str, Any],
        output_path: str
    ) -> Tuple[bool, str]:
        """
        Export report data to PowerPoint file

        Args:
            report_data: Report data dictionary with metadata and sections
            output_path: Path where PowerPoint file should be saved

        Returns:
            Tuple of (success: bool, result: str)
            - If success=True, result is the output file path
            - If success=False, result is the error message
        """
        try:
            # Validate report data structure
            validation_result = self._validate_report_data(report_data)
            if not validation_result[0]:
                return False, validation_result[1]

            # Ensure output directory exists
            output_path = Path(output_path).resolve()
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # Choose export method based on available dependencies
            if self.has_python_pptx:
                return self._export_with_pptx(report_data, output_path)
            else:
                return self._export_pdf_fallback(report_data, output_path)

        except Exception as e:
            return False, f"PowerPoint export failed: {str(e)}"

    def _validate_report_data(self, report_data: Dict[str, Any]) -> Tuple[bool, str]:
        """Validate report data structure"""
        if not isinstance(report_data, dict):
            return False, "report_data must be a dictionary"

        if "metadata" not in report_data:
            return False, "report_data must contain 'metadata' key"

        if "sections" not in report_data:
            return False, "report_data must contain 'sections' key"

        if not isinstance(report_data["sections"], list):
            return False, "'sections' must be a list"

        return True, "Valid"

    def _export_with_pptx(
        self,
        report_data: Dict[str, Any],
        output_path: Path
    ) -> Tuple[bool, str]:
        """Export PowerPoint using python-pptx (full-featured)"""
        try:
            # Create presentation
            prs = self.create_presentation()

            # Extract metadata
            metadata = report_data.get("metadata", {})
            report_type = metadata.get("report_type", "Analytics Report").title()
            generated_at = metadata.get("generated_at", datetime.now().isoformat())
            date_range = metadata.get("date_range", {})

            # Build title and subtitle
            title = f"{report_type} Report"
            subtitle_parts = []

            if date_range:
                start = date_range.get("start", "N/A")
                end = date_range.get("end", "N/A")
                subtitle_parts.append(f"Period: {start} to {end}")

            subtitle_parts.append(f"Generated: {self._format_datetime(generated_at)}")
            subtitle = "\n".join(subtitle_parts)

            # Add title slide
            self.add_title_slide(prs, title, subtitle)

            # Add table of contents
            sections = report_data.get("sections", [])
            if sections:
                toc_items = [section.get("title", "Section") for section in sections]
                self.add_text_slide(prs, "Table of Contents", toc_items)

            # Add section slides
            for section in sections:
                self._add_section_slide(prs, section)

            # Add closing slide
            self.add_text_slide(
                prs,
                "Questions?",
                ["Generated by dream-studio Analytics Platform", "", "Thank you for your attention"]
            )

            # Save presentation
            prs.save(str(output_path))

            return True, str(output_path)

        except Exception as e:
            return False, f"python-pptx export failed: {str(e)}"

    def create_presentation(self) -> Any:
        """
        Initialize new presentation with 16:9 layout

        Returns:
            Presentation object
        """
        if not self.has_python_pptx:
            raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")

        prs = Presentation()
        prs.slide_width = Inches(10)   # 16:9 aspect ratio
        prs.slide_height = Inches(5.625)

        return prs

    def add_title_slide(self, prs: Any, title: str, subtitle: str) -> None:
        """
        Add title slide to presentation

        Args:
            prs: Presentation object
            title: Main title text
            subtitle: Subtitle text
        """
        if not self.has_python_pptx:
            return

        # Use title slide layout (layout 0)
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = COLORS['primary']

        # Set subtitle
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            for paragraph in subtitle_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = COLORS['dark_gray']

    def add_chart_slide(self, prs: Any, chart_data: Dict[str, Any]) -> None:
        """
        Add slide with chart image

        Args:
            prs: Presentation object
            chart_data: Chart configuration with 'title', 'type', 'image_path' keys
        """
        if not self.has_python_pptx:
            return

        # Use title and content layout (layout 5)
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = chart_data.get("title", "Chart")
        slide.shapes.title.text = title
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = COLORS['primary']

        # Add chart image if available
        image_path = chart_data.get("image_path")
        if image_path and os.path.exists(image_path):
            try:
                # Position chart in center of content area
                left = Inches(1.5)
                top = Inches(2)
                width = Inches(7)
                height = Inches(3.5)

                slide.shapes.add_picture(image_path, left, top, width=width, height=height)
            except Exception as e:
                # Add placeholder text if image fails
                self._add_placeholder_text(slide, f"[Chart: {title} - Image error: {str(e)}]")
        else:
            # Add placeholder text
            chart_type = chart_data.get("type", "unknown")
            self._add_placeholder_text(slide, f"[Chart: {title} ({chart_type}) - Image not available]")

    def add_table_slide(self, prs: Any, table_data: Dict[str, Any]) -> None:
        """
        Add slide with data table

        Args:
            prs: Presentation object
            table_data: Dictionary with 'title' and 'data' keys
                       'data' should be dict of metric_name: value pairs
        """
        if not self.has_python_pptx:
            return

        # Use title and content layout
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = table_data.get("title", "Metrics")
        slide.shapes.title.text = title
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = COLORS['primary']

        # Prepare table data
        data = table_data.get("data", {})
        if not data:
            self._add_placeholder_text(slide, "[No data available]")
            return

        # Convert dict to rows
        rows = []
        rows.append(["Metric", "Value"])  # Header row

        for metric_name, metric_value in data.items():
            # Format value
            if isinstance(metric_value, (int, float)):
                value_str = f"{metric_value:,.2f}" if isinstance(metric_value, float) else f"{metric_value:,}"
            else:
                value_str = str(metric_value)

            rows.append([str(metric_name), value_str])

        # Create table
        num_rows = len(rows)
        num_cols = 2

        # Position table
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        row_height = Inches(0.4)
        height = row_height * num_rows

        # Add table shape
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # Populate table
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_value

                # Style cells
                text_frame = cell.text_frame
                paragraph = text_frame.paragraphs[0]

                if row_idx == 0:  # Header row
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(14)
                    paragraph.font.color.rgb = COLORS['white']
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['primary']
                else:  # Data rows
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = COLORS['primary']

                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = COLORS['light_gray']

                # Align columns
                if col_idx == 1:  # Value column - right align
                    paragraph.alignment = PP_ALIGN.RIGHT
                else:  # Metric column - left align
                    paragraph.alignment = PP_ALIGN.LEFT

        # Set column widths
        table.columns[0].width = Inches(4.5)
        table.columns[1].width = Inches(2.5)

    def add_text_slide(self, prs: Any, title: str, content: List[str]) -> None:
        """
        Add text/bullet point slide

        Args:
            prs: Presentation object
            title: Slide title
            content: List of text items (bullet points)
        """
        if not self.has_python_pptx:
            return

        # Use title and content layout
        slide_layout = prs.slide_layouts[1]  # Bullet layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        slide.shapes.title.text = title
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = COLORS['primary']

        # Add content
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            for i, item in enumerate(content):
                if i == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()

                paragraph.text = item
                paragraph.level = 0
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = COLORS['primary']

    def _add_section_slide(self, prs: Any, section: Dict[str, Any]) -> None:
        """
        Add slide(s) for a report section

        Args:
            prs: Presentation object
            section: Section data with title, metrics, charts
        """
        section_title = section.get("title", "Section")
        metrics = section.get("metrics", {})
        charts = section.get("charts", [])

        # If section has metrics, create table slide
        if metrics:
            self.add_table_slide(prs, {"title": section_title, "data": metrics})

        # If section has charts, create chart slide(s)
        for chart in charts:
            chart_with_section = chart.copy()
            if not chart_with_section.get("title"):
                chart_with_section["title"] = f"{section_title} - Chart"

            self.add_chart_slide(prs, chart_with_section)

        # If section has neither metrics nor charts, create text slide
        if not metrics and not charts:
            self.add_text_slide(prs, section_title, ["No data available for this section"])

    def _add_placeholder_text(self, slide: Any, text: str) -> None:
        """Add placeholder text to a slide"""
        left = Inches(2)
        top = Inches(2.5)
        width = Inches(6)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text

        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(16)
        paragraph.font.italic = True
        paragraph.font.color.rgb = COLORS['dark_gray']
        paragraph.alignment = PP_ALIGN.CENTER

    def _export_pdf_fallback(
        self,
        report_data: Dict[str, Any],
        output_path: Path
    ) -> Tuple[bool, str]:
        """Export as PDF when python-pptx unavailable"""
        try:
            from analytics.exporters.pdf_exporter import PDFExporter

            # Change extension to .pdf
            pdf_path = output_path.with_suffix(".pdf")

            # Export using PDF exporter
            pdf_exporter = PDFExporter()
            success, result = pdf_exporter.export_to_pdf(report_data, str(pdf_path))

            if success:
                # Return the actual file path, not a message
                # The user will see the warning message at initialization
                return True, result
            else:
                return False, result

        except ImportError:
            return False, (
                "python-pptx not available and PDF fallback failed. "
                "Install dependencies with: pip install python-pptx reportlab Pillow"
            )
        except Exception as e:
            return False, f"PDF fallback export failed: {str(e)}"

    def _format_datetime(self, dt_str: str) -> str:
        """Format ISO datetime string for display"""
        try:
            dt = datetime.fromisoformat(dt_str.replace('Z', '+00:00'))
            return dt.strftime("%Y-%m-%d %H:%M:%S UTC")
        except Exception:
            return dt_str

    def get_installation_instructions(self) -> str:
        """Get installation instructions for PowerPoint libraries"""
        if self.has_python_pptx:
            return "python-pptx is installed and ready to use."
        else:
            return (
                "PowerPoint export requires python-pptx. Install with:\n"
                "  pip install python-pptx Pillow\n"
                "\n"
                "PDF export is available as fallback."
            )


def main():
    """Demo usage of PPTXExporter"""
    from datetime import datetime

    # Sample report data
    sample_data = {
        "metadata": {
            "generated_at": datetime.now().isoformat(),
            "report_type": "executive",
            "date_range": {"start": "2026-04-01", "end": "2026-04-30"}
        },
        "sections": [
            {
                "title": "Key Performance Indicators",
                "metrics": {
                    "Total Users": 15234,
                    "Active Users": 12458,
                    "Conversion Rate": 0.0823,
                    "Revenue": 125430.50,
                    "Churn Rate": 0.0234
                },
                "charts": []
            },
            {
                "title": "User Growth Trends",
                "metrics": {},
                "charts": [
                    {
                        "title": "Monthly Active Users",
                        "type": "line",
                        "image_path": None  # Would be path to chart image
                    }
                ]
            },
            {
                "title": "Revenue Analysis",
                "metrics": {
                    "Q1 Revenue": 125430.50,
                    "Growth Rate": 0.15,
                    "Target Achievement": 0.95
                },
                "charts": []
            }
        ]
    }

    # Create exporter
    exporter = PPTXExporter()

    # Show library info
    print(exporter.get_installation_instructions())
    print()

    # Export presentation
    output_path = "C:/Users/Dannis Seay/Downloads/sample_analytics_report.pptx"
    success, result = exporter.export_to_pptx(sample_data, output_path)

    if success:
        print(f"Presentation exported successfully")
        print(f"  Path: {result}")
    else:
        print(f"Export failed: {result}")


if __name__ == "__main__":
    main()

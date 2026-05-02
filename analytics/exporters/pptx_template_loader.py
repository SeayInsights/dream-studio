"""
PowerPoint Template Loader

Loads and applies professional PowerPoint templates for dream-studio analytics exports.
"""

import os
from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class PPTXTemplateLoader:
    """Helper class to load and apply PowerPoint templates."""

    def __init__(self):
        """Initialize the template loader."""
        self.template_dir = Path(__file__).parent / "pptx_templates"
        self.template_dir.mkdir(exist_ok=True)

        # Color scheme
        self.colors = {
            'primary': RGBColor(44, 62, 80),      # Dark blue #2c3e50
            'secondary': RGBColor(52, 152, 219),  # Bright blue #3498db
            'accent': RGBColor(230, 126, 34),     # Orange accent #e67e22
            'success': RGBColor(39, 174, 96),     # Green #27ae60
            'warning': RGBColor(241, 196, 15),    # Yellow #f1c40f
            'danger': RGBColor(231, 76, 60),      # Red #e74c3c
            'light': RGBColor(236, 240, 241),     # Light gray #ecf0f1
            'dark': RGBColor(52, 73, 94),         # Dark gray #34495e
        }

    def load_template(self, template_name: str) -> Presentation:
        """
        Load a template presentation by name.

        Args:
            template_name: Name of template ('executive' or 'technical')

        Returns:
            Presentation object loaded from template

        Raises:
            FileNotFoundError: If template doesn't exist
        """
        template_path = self.template_dir / f"{template_name}.pptx"

        if not template_path.exists():
            raise FileNotFoundError(
                f"Template '{template_name}' not found at {template_path}. "
                f"Available templates: {', '.join(self.list_available_templates())}"
            )

        return Presentation(str(template_path))

    def apply_template(
        self,
        prs: Presentation,
        template_name: str
    ) -> Presentation:
        """
        Apply template styling to an existing presentation.

        Args:
            prs: Presentation to apply template to
            template_name: Name of template to apply

        Returns:
            Modified presentation with template styling applied
        """
        template = self.load_template(template_name)

        # Copy slide master from template
        if template.slide_master:
            # Apply color scheme and fonts from template
            for slide in prs.slides:
                # Apply background
                if hasattr(slide, 'background'):
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        return prs

    def list_available_templates(self) -> List[str]:
        """
        List all available template names.

        Returns:
            List of template names (without .pptx extension)
        """
        templates = []
        if self.template_dir.exists():
            for file in self.template_dir.glob("*.pptx"):
                if not file.name.startswith("~"):  # Skip temp files
                    templates.append(file.stem)
        return sorted(templates)

    def create_executive_template(self) -> Presentation:
        """
        Create executive summary template programmatically.

        Returns:
            Presentation object with executive template structure
        """
        prs = Presentation()
        prs.slide_width = Inches(10)   # 16:9 aspect ratio
        prs.slide_height = Inches(5.625)

        # Slide 1: Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_title_slide(slide, "Dream-Studio Analytics Report",
                             "[Report Period]")

        # Slide 2: Executive Summary
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_content_slide(slide, "Executive Summary",
                               ["Key finding 1", "Key finding 2", "Key finding 3",
                                "Key finding 4", "Key finding 5"])

        # Slide 3: Key Metrics Dashboard
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_metrics_slide(slide, "Key Metrics")

        # Slide 4: Trends & Insights
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_chart_slide(slide, "Trends & Insights")

        # Slide 5: Recommendations
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_content_slide(slide, "Recommendations",
                               ["Priority 1: [Recommendation]",
                                "Priority 2: [Recommendation]",
                                "Priority 3: [Recommendation]"])

        # Slide 6: Q&A
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_qa_slide(slide)

        return prs

    def create_technical_template(self) -> Presentation:
        """
        Create technical deep-dive template programmatically.

        Returns:
            Presentation object with technical template structure
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        slides_content = [
            ("Title", "title", None),
            ("Table of Contents", "toc", None),
            ("Methodology & Data Sources", "content",
             ["Data collection period", "Analysis framework", "Data sources"]),
            ("Skill Usage Overview", "chart", None),
            ("Skill Performance Deep-Dive", "chart", None),
            ("Token Usage Analysis", "chart", None),
            ("Cost Analysis", "metrics", None),
            ("Session Analytics", "chart", None),
            ("Session Duration Patterns", "chart", None),
            ("Model Performance Comparison", "chart", None),
            ("Model Selection Patterns", "content",
             ["Usage by model", "Cost per model", "Performance metrics"]),
            ("Lessons Learned", "content",
             ["Key insight 1", "Key insight 2", "Key insight 3"]),
            ("Workflow Analysis", "chart", None),
            ("Workflow Efficiency Metrics", "metrics", None),
            ("Detailed Recommendations", "content",
             ["Recommendation 1", "Recommendation 2", "Recommendation 3"]),
            ("Appendix: Methodology Details", "content",
             ["Data processing", "Analysis approach", "Limitations"]),
            ("Contact & Questions", "qa", None),
        ]

        for idx, (title, slide_type, content) in enumerate(slides_content, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            if slide_type == "title":
                self._add_title_slide(slide,
                                     "Dream-Studio Technical Analytics",
                                     "[Detailed Analysis Report]")
            elif slide_type == "toc":
                self._add_toc_slide(slide)
            elif slide_type == "content":
                self._add_content_slide(slide, title, content or [])
            elif slide_type == "chart":
                self._add_chart_slide(slide, title)
            elif slide_type == "metrics":
                self._add_metrics_slide(slide, title)
            elif slide_type == "qa":
                self._add_qa_slide(slide)

            # Add slide number to footer
            self._add_footer(slide, idx, len(slides_content))

        return prs

    def _add_title_slide(self, slide, title: str, subtitle: str):
        """Add title slide with branding."""
        # Background
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.8), Inches(8), Inches(0.6)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['secondary']

        # Accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(3.6), Inches(3), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.fill.background()

    def _add_content_slide(self, slide, title: str, bullets: List[str]):
        """Add content slide with title and bullet points."""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # Accent line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.95), Inches(2), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['secondary']
        line.line.fill.background()

        # Content area
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark']
            p.level = 0

    def _add_chart_slide(self, slide, title: str):
        """Add slide with chart placeholder."""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # Chart placeholder
        chart_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.8)
        )
        chart_placeholder.fill.solid()
        chart_placeholder.fill.fore_color.rgb = self.colors['light']
        chart_placeholder.line.color.rgb = self.colors['secondary']
        chart_placeholder.line.width = Pt(2)

        # Placeholder text
        text_frame = chart_placeholder.text_frame
        text_frame.text = "[CHART PLACEHOLDER]"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors['dark']

        # Insights box
        insights_box = slide.shapes.add_textbox(
            Inches(6.5), Inches(1.3), Inches(3), Inches(3.8)
        )
        insights_frame = insights_box.text_frame
        insights_frame.word_wrap = True

        # Insights title
        p = insights_frame.paragraphs[0]
        p.text = "Key Insights"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # Bullet points
        for i in range(3):
            p = insights_frame.add_paragraph()
            p.text = f"Insight {i+1}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['dark']
            p.level = 0

    def _add_metrics_slide(self, slide, title: str):
        """Add slide with metric cards."""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # Metric cards (2x2 grid)
        metrics = [
            ("Total Sessions", "[Value]", self.colors['secondary']),
            ("Total Tokens", "[Value]", self.colors['success']),
            ("Total Cost", "[Value]", self.colors['warning']),
            ("Avg Session", "[Value]", self.colors['accent']),
        ]

        positions = [
            (0.8, 1.5), (5.3, 1.5),
            (0.8, 3.5), (5.3, 3.5)
        ]

        for (metric_title, value, color), (x, y) in zip(metrics, positions):
            self._add_metric_card(slide, x, y, metric_title, value, color)

    def _add_metric_card(self, slide, x: float, y: float,
                        title: str, value: str, color: RGBColor):
        """Add a single metric card."""
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(4), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # Metric value
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 0.2), Inches(3.4), Inches(0.6)
        )
        value_frame = value_box.text_frame
        value_frame.text = value
        p = value_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color

        # Metric title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 0.9), Inches(3.4), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['dark']

    def _add_toc_slide(self, slide):
        """Add table of contents slide."""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Table of Contents"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # TOC sections
        sections = [
            "Methodology & Data Sources",
            "Skill Usage Analysis",
            "Token & Cost Analysis",
            "Session Analytics",
            "Model Performance",
            "Lessons Learned",
            "Workflow Analysis",
            "Recommendations",
            "Appendix"
        ]

        # Two column layout
        left_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.3), Inches(4), Inches(3.8)
        )
        right_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.3), Inches(4), Inches(3.8)
        )

        mid = len(sections) // 2

        for i, section in enumerate(sections[:mid]):
            p = left_box.text_frame.add_paragraph() if i > 0 else left_box.text_frame.paragraphs[0]
            p.text = f"{i+1}. {section}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark']

        for i, section in enumerate(sections[mid:], mid):
            p = right_box.text_frame.add_paragraph() if i > mid else right_box.text_frame.paragraphs[0]
            p.text = f"{i+1}. {section}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark']

    def _add_qa_slide(self, slide):
        """Add Q&A / Contact slide."""
        # Background
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['primary']

        # Q&A text
        qa_box = slide.shapes.add_textbox(
            Inches(2), Inches(1.5), Inches(6), Inches(1.5)
        )
        qa_frame = qa_box.text_frame
        qa_frame.text = "Questions?"
        p = qa_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5), Inches(6), Inches(1)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "dream-studio analytics\ndannis.seay@twinrootsllc.com"
        for p in contact_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['secondary']

    def _add_footer(self, slide, slide_num: int, total_slides: int):
        """Add footer with slide number."""
        footer_box = slide.shapes.add_textbox(
            Inches(9), Inches(5.2), Inches(0.8), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = f"{slide_num}/{total_slides}"
        p = footer_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['dark']


def create_templates():
    """
    Create both executive and technical templates.

    Run this function to generate the template files.
    """
    loader = PPTXTemplateLoader()

    # Create executive template
    print("Creating executive template...")
    exec_prs = loader.create_executive_template()
    exec_path = loader.template_dir / "executive.pptx"
    exec_prs.save(str(exec_path))
    print(f"[OK] Created: {exec_path}")

    # Create technical template
    print("Creating technical template...")
    tech_prs = loader.create_technical_template()
    tech_path = loader.template_dir / "technical.pptx"
    tech_prs.save(str(tech_path))
    print(f"[OK] Created: {tech_path}")

    print(f"\nTemplates created successfully!")
    print(f"Available templates: {', '.join(loader.list_available_templates())}")


if __name__ == "__main__":
    create_templates()

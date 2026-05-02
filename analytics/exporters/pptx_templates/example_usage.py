"""
Example Usage: PowerPoint Template Loader

Demonstrates how to use the PPTX templates for dream-studio analytics reports.
"""

import sys
from pathlib import Path

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent))

from analytics.exporters.pptx_template_loader import PPTXTemplateLoader
from pptx.util import Inches, Pt


def example_executive_report():
    """Example: Create an executive report from template."""
    print("Creating executive report...")

    loader = PPTXTemplateLoader()

    # Load executive template
    prs = loader.load_template("executive")

    # Slide 1: Title (already formatted)
    # Slide 2: Executive Summary
    slide = prs.slides[1]
    # In production, you would replace placeholder text with actual data
    # For this example, we'll just note what would be replaced
    print("  - Title slide ready")
    print("  - Executive summary slide ready (5 bullet points)")

    # Slide 3: Key Metrics Dashboard
    # In production: update the 4 metric cards with actual values
    print("  - Key metrics slide ready (4 metric cards)")

    # Slide 4: Trends & Insights
    # In production: insert actual chart and update insights
    print("  - Trends slide ready (chart placeholder + insights)")

    # Slide 5: Recommendations
    print("  - Recommendations slide ready")

    # Slide 6: Q&A
    print("  - Q&A slide ready")

    # Save
    output_path = Path(__file__).parent / "example_executive_output.pptx"
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def example_technical_report():
    """Example: Create a technical report from template."""
    print("Creating technical report...")

    loader = PPTXTemplateLoader()

    # Load technical template
    prs = loader.load_template("technical")

    print(f"  - Loaded template with {len(prs.slides)} slides")

    # In production, you would:
    # 1. Load analytics data from SQLite database
    # 2. Generate charts using matplotlib/plotly
    # 3. Insert charts into chart placeholder areas
    # 4. Update text boxes with actual data
    # 5. Populate metrics cards with real values

    print("  - All slides ready for data population")

    # Save
    output_path = Path(__file__).parent / "example_technical_output.pptx"
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def example_list_templates():
    """Example: List available templates."""
    print("Available templates:")

    loader = PPTXTemplateLoader()
    templates = loader.list_available_templates()

    for template in templates:
        print(f"  - {template}")

    return templates


def example_custom_slide():
    """Example: Add a custom slide to an existing template."""
    print("Creating custom slide example...")

    loader = PPTXTemplateLoader()
    prs = loader.load_template("executive")

    # Add a new custom slide after slide 3
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add custom title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Custom Slide: Additional Analysis"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = loader.colors['primary']

    # Add custom content
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8)
    )
    text_frame = content_box.text_frame
    text_frame.text = "This is a custom slide added programmatically.\n\n"

    p = text_frame.add_paragraph()
    p.text = "You can add any content you need:"
    p.font.size = Pt(18)

    for item in ["Custom charts", "Additional metrics", "Supplementary analysis"]:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 1

    # Save
    output_path = Path(__file__).parent / "example_custom_slide.pptx"
    prs.save(str(output_path))
    print(f"Saved: {output_path}")


def main():
    """Run all examples."""
    print("=" * 60)
    print("PowerPoint Template Loader - Usage Examples")
    print("=" * 60)
    print()

    # Example 1: List templates
    print("Example 1: List Available Templates")
    print("-" * 60)
    templates = example_list_templates()
    print()

    # Example 2: Executive report
    print("Example 2: Executive Report")
    print("-" * 60)
    example_executive_report()
    print()

    # Example 3: Technical report
    print("Example 3: Technical Report")
    print("-" * 60)
    example_technical_report()
    print()

    # Example 4: Custom slide
    print("Example 4: Custom Slide")
    print("-" * 60)
    example_custom_slide()
    print()

    print("=" * 60)
    print("All examples completed!")
    print("=" * 60)
    print()
    print("Next steps:")
    print("1. Open the generated .pptx files to review")
    print("2. Integrate with PPTXExporter for automated reports")
    print("3. Customize templates by editing the .pptx files directly")
    print("4. Regenerate templates: py pptx_template_loader.py")


if __name__ == "__main__":
    main()
